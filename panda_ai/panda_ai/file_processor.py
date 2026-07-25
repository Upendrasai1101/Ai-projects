# --- FILE: file_processor.py ---
"""
file_processor.py — Panda AI V8.2 Cloud Multimodal Processor

ARCHITECTURE CHANGE (V8.2):
  OLD: Local libraries (whisper, pytesseract) — crashed localhost, broke HF Spaces
  NEW: Cloud API calls only — zero heavy dependencies, works on HF Spaces

  Audio/Video  → Groq Whisper API  (distil-whisper-large-v3-en)
                 Key passed in from app.py via set_groq_key()
                 No local model download. Streams bytes directly.

  Images       → HF Inference API  (Salesforce/blip-image-captioning-large)
                 HF_TOKEN env var. Returns rich image description.
                 Falls back to filename hint if token missing.

  PDF/DOCX/XLSX/PPTX — unchanged from V5 (text extraction, no model needed)

WHAT IS UNTOUCHED:
  process_file(), process_mixed_files(), allowed_file(), UPLOAD_DIR,
  chunk_text(), find_relevant_chunks(), clean_extracted_text(),
  extract_pdf(), extract_word(), extract_excel(), extract_pptx()
  — all identical to V5.
"""

import os
import re
import random
import requests

UPLOAD_DIR = "uploads"
os.makedirs(f"{UPLOAD_DIR}/documents", exist_ok=True)
os.makedirs(f"{UPLOAD_DIR}/images",    exist_ok=True)
os.makedirs(f"{UPLOAD_DIR}/audio",     exist_ok=True)

ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'xlsx', 'xls', 'pptx',
    'jpg', 'jpeg', 'png', 'bmp', 'tiff',
    'mp3', 'wav', 'ogg', 'mp4', 'avi', 'mov'
}

# ── V8.3: Dedicated file processor key pools ─────────────────────────────
# Completely separate from chat pool (GROQ_API_KEY_1..8 in app.py)
# app.py no longer injects a key here — file_processor reads env directly.

# Vision pool — image analysis
_VISION_KEYS = [k for k in [
    os.getenv("GROQ_VISION_KEY_1"),
] if k]

# Whisper pool — audio + video transcription
_WHISPER_KEYS = [k for k in [
    os.getenv("GROQ_WHISPER_KEY_1"),
] if k]

# Model names — updated (deprecated models removed)
GROQ_VISION_MODEL   = "qwen/qwen3.6-27b"
GROQ_WHISPER_MODEL  = "whisper-large-v3-turbo"
GROQ_CHAT_ENDPOINT  = "https://api.groq.com/openai/v1/chat/completions"
GROQ_AUDIO_ENDPOINT = "https://api.groq.com/openai/v1/audio/transcriptions"

# HF token — kept as optional fallback only (not primary anymore)
_hf_token = os.getenv("HF_TOKEN", "")


def _get_vision_key():
    """Pick a vision key. Returns None if not configured."""
    return random.choice(_VISION_KEYS) if _VISION_KEYS else None


def _get_whisper_key():
    """Pick a whisper key. Returns None if not configured."""
    return random.choice(_WHISPER_KEYS) if _WHISPER_KEYS else None

# ── UNCHANGED HELPERS ───────────────────────────────────────────────────────
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_extension(filename):
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

def clean_extracted_text(text, limit=8000):
    if not text:
        return ""
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    text = text.strip()
    return text[:limit]

def chunk_text(text, chunk_size=1500):
    """Split text into chunks of ~chunk_size words."""
    words  = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

def find_relevant_chunks(chunks, question, top_n=2):
    """Find most relevant chunks based on keyword overlap."""
    if not chunks:
        return ""
    q_words = set(re.sub(r'[^\w\s]', '', question.lower()).split())
    scores  = []
    for chunk in chunks:
        c_words = set(chunk.lower().split())
        score   = len(q_words & c_words)
        scores.append(score)
    sorted_idx = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)
    top_chunks = [chunks[i] for i in sorted_idx[:top_n]]
    return "\n\n---\n\n".join(top_chunks)

# ════════════════════════════════════════
# PDF Extraction  (UNCHANGED)
# ════════════════════════════════════════
def extract_pdf(file_path):
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        full_text = "\n\n".join(text_parts)
        cleaned_text = clean_extracted_text(full_text)
        filename = os.path.basename(file_path)
        formatted_text = (
            f"[DOCUMENT CONTENT ATTACHED BELOW]\n\n"
            f"The user uploaded a PDF named '{filename}'. Here is the extracted text:\n"
            f"---\n{cleaned_text}\n---\n"
            "Instruct the LLM to read the above text and directly answer the user's question."
        )
        print(f"📄 PDF extracted {len(cleaned_text)} characters")
        return formatted_text
    except ImportError:
        return "[PDF extraction unavailable: pdfplumber not installed]"
    except Exception as e:
        print(f"PDF error: {e}")
        return f"[PDF extraction failed: {e}]"

# ════════════════════════════════════════
# Word (.docx) Extraction  (UNCHANGED)
# ════════════════════════════════════════
def extract_word(file_path):
    try:
        from docx import Document
        doc   = Document(file_path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        full  = "\n".join(paras)
        print(f"DOCX extracted: {len(full)} chars")
        return clean_extracted_text(full)
    except ImportError:
        return "[Word extraction unavailable: python-docx not installed]"
    except Exception as e:
        print(f"DOCX error: {e}")
        return f"[Word extraction failed: {e}]"

# ════════════════════════════════════════
# Excel (.xlsx) Extraction  (UNCHANGED)
# ════════════════════════════════════════
def extract_excel(file_path):
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(file_path, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        full = "\n".join(parts)
        print(f"Excel extracted: {len(full)} chars")
        return clean_extracted_text(full)
    except ImportError:
        return "[Excel extraction unavailable: openpyxl not installed]"
    except Exception as e:
        print(f"Excel error: {e}")
        return f"[Excel extraction failed: {e}]"

# ════════════════════════════════════════
# PowerPoint (.pptx) Extraction  (UNCHANGED)
# ════════════════════════════════════════
def extract_pptx(file_path):
    try:
        from pptx import Presentation
        prs   = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        full = "\n".join(parts)
        print(f"PPTX extracted: {len(full)} chars")
        return clean_extracted_text(full)
    except ImportError:
        return "[PPT extraction unavailable: python-pptx not installed]"
    except Exception as e:
        print(f"PPTX error: {e}")
        return f"[PPT extraction failed: {e}]"

# ════════════════════════════════════════
# V8.3 — Image analysis via Groq Qwen Vision
# Replaces HF BLIP — HF_TOKEN is no longer required.
# ════════════════════════════════════════
def extract_image(file_path):
    """
    V8.3 — Image analysis via Groq Qwen Vision (qwen/qwen3.6-27b)
    Replaces HF BLIP — HF_TOKEN is not required.

    Returns: "[Image Analysis]: <description>"
    """
    print(f"🖼️  Image → Qwen Vision: {os.path.basename(file_path)}")

    key = _get_vision_key()
    if not key:
        return (
            f"[Image: {os.path.basename(file_path)} — "
            "GROQ_VISION_KEY_1 not set in .env / HF Spaces secrets. "
            "Add the vision account key to enable image analysis.]"
        )

    try:
        import base64
        with open(file_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()

        ext = file_path.rsplit('.', 1)[-1].lower()
        mime = {
            'jpg':  'image/jpeg', 'jpeg': 'image/jpeg',
            'png':  'image/png',  'bmp':  'image/bmp',
            'tiff': 'image/tiff',
        }.get(ext, 'image/jpeg')

        res = requests.post(
            GROQ_CHAT_ENDPOINT,
            headers={
                "Authorization": f"Bearer {key}",
                "Content-Type":  "application/json",
            },
            json={
                "model": GROQ_VISION_MODEL,
                "messages": [{
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime};base64,{b64}"}
                        },
                        {
                            "type": "text",
                            "text": (
                                "Describe this image in detail. "
                                "If it is a car, identify the make, model, and approximate year. "
                                "If it contains text, extract all visible text. "
                                "If it is a chart or table, explain the data. "
                                "If it is a document or screenshot, summarize the content. "
                                "Be specific and thorough."
                            )
                        }
                    ]
                }],
                "max_tokens": 500,
            },
            timeout=30,
        )

        if res.status_code == 429:
            print("  ⚠️  Qwen Vision rate limited")
            return "[Image analysis rate limited. Please try again shortly.]"

        if res.status_code != 200:
            print(f"  ❌ Qwen Vision error {res.status_code}: {res.text[:200]}")
            return f"[Image analysis failed: HTTP {res.status_code}]"

        desc = res.json()["choices"][0]["message"]["content"].strip()
        print(f"  ✅ Qwen Vision OK: {len(desc)} chars")
        return f"[Image Analysis]: {desc}"

    except requests.exceptions.Timeout:
        print("  ❌ Qwen Vision timeout")
        return "[Image analysis timed out. Please try again.]"
    except Exception as e:
        print(f"  ❌ Image extraction error: {e}")
        return f"[Image extraction failed: {e}]"

# ════════════════════════════════════════
# V8.2 — Audio Transcription via Groq Whisper API
# Replaces: local openai-whisper (1.5 GB model download)
# Model: distil-whisper-large-v3-en
# ════════════════════════════════════════
def extract_audio(file_path):
    """
    Transcribe audio by sending bytes to Groq's Whisper endpoint.

    Why Groq Whisper over local whisper?
      - Local whisper downloads a 1.5 GB model → crashes HF Spaces RAM limit.
      - Groq transcribes in ~2-5 seconds via API with no local storage.
      - distil-whisper-large-v3-en = same quality, 6x faster than large-v3.

    Supported formats: mp3, wav, ogg (Groq accepts all standard audio).
    File size limit: 25 MB (Groq limit). Larger files are rejected gracefully.

    Fallback: If Groq key is unavailable, returns a clear error message
    so the AI can ask the user to describe the audio content manually.
    """
    print(f"🎵  Audio → Groq Whisper: {os.path.basename(file_path)}")

    whisper_key = _get_whisper_key()
    if not whisper_key:
        print("  ⚠️  Groq whisper key not available for audio transcription")
        return (
            "[Audio: GROQ_WHISPER_KEY_1 not set. "
            "Add the whisper account key to .env / HF Spaces secrets.]"
        )

    # Check file size — Groq limit is 25 MB
    file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
    if file_size_mb > 25:
        return (
            f"[Audio file too large ({file_size_mb:.1f} MB). "
            "Groq Whisper accepts files up to 25 MB. "
            "Please trim the audio or describe its content.]"
        )

    try:
        filename = os.path.basename(file_path)
        ext      = filename.rsplit('.', 1)[-1].lower()

        # Map extensions to MIME types for the multipart upload
        mime_map = {
            'mp3': 'audio/mpeg',
            'wav': 'audio/wav',
            'ogg': 'audio/ogg',
            'm4a': 'audio/mp4',
            'flac': 'audio/flac',
            'webm': 'audio/webm',
        }
        mime_type = mime_map.get(ext, 'audio/mpeg')

        with open(file_path, "rb") as f:
            audio_bytes = f.read()

        response = requests.post(
            GROQ_AUDIO_ENDPOINT,
            headers={"Authorization": f"Bearer {whisper_key}"},
            files={
                "file": (filename, audio_bytes, mime_type),
            },
            data={
                "model":           GROQ_WHISPER_MODEL,
                "response_format": "json",
                "language":        "en",   # set to None for auto-detect if needed
            },
            timeout=60,
        )

        if response.status_code == 429:
            print("  ⚠️  Groq Whisper rate limited")
            return "[Audio transcription rate limited. Please try again in a moment.]"

        if response.status_code != 200:
            print(f"  ❌ Groq Whisper error {response.status_code}: {response.text[:200]}")
            return f"[Audio transcription failed (HTTP {response.status_code}).]"

        result = response.json()
        text   = result.get("text", "").strip()

        if not text:
            return "[No speech detected in the audio file.]"

        print(f"  ✅ Transcribed: {len(text)} chars")
        return clean_extracted_text(text)

    except requests.exceptions.Timeout:
        print("  ❌ Groq Whisper timeout")
        return "[Audio transcription timed out. Please try a shorter clip.]"
    except Exception as e:
        print(f"  ❌ Audio extraction error: {e}")
        return f"[Audio extraction failed: {e}]"

# ════════════════════════════════════════
# V8.2 — Video Processing
# Replaces: moviepy (heavy) + local whisper
# Strategy: extract audio track → Groq Whisper
# Uses ffmpeg (available on HF Spaces by default)
# ════════════════════════════════════════
def extract_video(file_path):
    """
    Extract audio from video using ffmpeg (pre-installed on HF Spaces),
    then transcribe via Groq Whisper.

    Why ffmpeg instead of moviepy?
      - moviepy pulls in numpy, imageio, etc. — heavy and slow to install.
      - ffmpeg is already present in HF Spaces Docker image.
      - subprocess call is instant — no Python import overhead.

    If ffmpeg is unavailable, falls back to sending the raw video bytes
    to Groq Whisper directly (Groq accepts mp4/webm natively).
    """
    print(f"🎬  Video → ffmpeg extract → Groq Whisper: {os.path.basename(file_path)}")

    audio_path = os.path.join(UPLOAD_DIR, "audio", "temp_video_audio.wav")

    # Strategy 1: extract audio with ffmpeg
    try:
        import subprocess
        result = subprocess.run(
            [
                "ffmpeg", "-y",          # overwrite without prompt
                "-i", file_path,         # input video
                "-vn",                   # no video stream
                "-acodec", "pcm_s16le",  # raw WAV
                "-ar", "16000",          # 16kHz — Whisper optimal sample rate
                "-ac", "1",              # mono
                audio_path,
            ],
            capture_output=True,
            timeout=60,
        )

        if result.returncode == 0 and os.path.exists(audio_path):
            print(f"  ✅ ffmpeg extracted audio → {audio_path}")
            transcript = extract_audio(audio_path)
            try: os.remove(audio_path)
            except: pass
            return transcript
        else:
            err = result.stderr.decode("utf-8", errors="ignore")[:200]
            print(f"  ⚠️  ffmpeg failed: {err}")

    except FileNotFoundError:
        print("  ⚠️  ffmpeg not found — trying direct Groq upload")
    except Exception as e:
        print(f"  ⚠️  ffmpeg error: {e}")

    # Strategy 2: send video directly to Groq (supports mp4/webm)
    print("  🔄 Fallback: sending video bytes directly to Groq Whisper")
    try:
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > 25:
            return (
                f"[Video file too large ({file_size_mb:.1f} MB) for direct transcription. "
                "Use a shorter clip under 25 MB.]"
            )

        whisper_key = _get_whisper_key()
        if not whisper_key:
            return "[Video: GROQ_WHISPER_KEY_1 not set.]"

        filename = os.path.basename(file_path)
        ext      = filename.rsplit('.', 1)[-1].lower()
        mime_map = {'mp4':'video/mp4','webm':'video/webm','mov':'video/quicktime','avi':'video/x-msvideo'}
        mime_type = mime_map.get(ext, 'video/mp4')

        with open(file_path, "rb") as f:
            video_bytes = f.read()

        response = requests.post(
            GROQ_AUDIO_ENDPOINT,
            headers={"Authorization": f"Bearer {whisper_key}"},
            files={"file": (filename, video_bytes, mime_type)},
            data={"model": GROQ_WHISPER_MODEL, "response_format": "json"},
            timeout=90,
        )

        if response.status_code == 200:
            text = response.json().get("text", "").strip()
            if text:
                print(f"  ✅ Direct video transcription: {len(text)} chars")
                return clean_extracted_text(text)
            return "[No speech detected in the video.]"

        print(f"  ❌ Direct Groq video failed: {response.status_code}")
        return f"[Video transcription failed (HTTP {response.status_code}).]"

    except Exception as e:
        print(f"  ❌ Video fallback error: {e}")
        return f"[Video extraction failed: {e}]"

# ════════════════════════════════════════
# MAIN: process_file()  (UNCHANGED signature)
# ════════════════════════════════════════
def process_file(file_path, filename):
    """
    Main entry point — detects file type and extracts text/description.
    Returns: (extracted_text, file_type)
    Signature identical to V5 — app.py needs zero changes here.
    """
    ext = get_extension(filename)

    if ext == 'pdf':
        return extract_pdf(file_path),   'PDF'
    elif ext == 'docx':
        return extract_word(file_path),  'Word'
    elif ext in ('xlsx', 'xls'):
        return extract_excel(file_path), 'Excel'
    elif ext == 'pptx':
        return extract_pptx(file_path),  'PowerPoint'
    elif ext in ('jpg', 'jpeg', 'png', 'bmp', 'tiff'):
        return extract_image(file_path), 'Image'
    elif ext in ('mp3', 'wav', 'ogg'):
        return extract_audio(file_path), 'Audio'
    elif ext in ('mp4', 'avi', 'mov'):
        return extract_video(file_path), 'Video'
    else:
        return "[Unsupported file type]", 'Unknown'

# ════════════════════════════════════════
# Mixed Input  (UNCHANGED)
# ════════════════════════════════════════
def process_mixed_files(file_paths_names, question):
    """
    Handle multiple files uploaded together.
    Returns combined relevant context for the AI.
    """
    all_sections = []

    for file_path, filename in file_paths_names:
        text, ftype = process_file(file_path, filename)
        if text and not text.startswith('['):
            all_sections.append(f"=== {ftype} File: {filename} ===\n{text}")

    if not all_sections:
        return ""

    combined = "\n\n".join(all_sections)
    chunks   = chunk_text(combined, chunk_size=1500)
    relevant = find_relevant_chunks(chunks, question, top_n=2)

    print(f"Mixed input: {len(all_sections)} files, {len(chunks)} chunks, relevant: {len(relevant)} chars")
    return relevant